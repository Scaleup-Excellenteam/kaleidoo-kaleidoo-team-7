import pytesseract
from PIL import Image, ImageEnhance, ImageFilter
import pandas as pd
import os
import glob
import re
from langdetect import detect_langs, DetectorFactory
from pptx import Presentation
from docx import Document
from pdfminer.high_level import extract_text
from PyPDF2 import PdfReader
import arabic_reshaper
from docx.enum.text import WD_BREAK
from typing import Optional, List
import fitz  # PyMuPDF for PDF processing
from pprint import pprint
from TextProcessorSingleton import TextProcessorSingleton


text_processor = TextProcessorSingleton.get_instance()

# Ensure consistent language detection
DetectorFactory.seed = 0

# Extended language map (detect to Tesseract)
lang_map = {
    'en': 'eng',    # English
    'es': 'spa',    # Spanish
    'fr': 'fra',    # French
    'de': 'deu',    # German
    'it': 'ita',    # Italian
    'pt': 'por',    # Portuguese
    'ru': 'rus',    # Russian
    'zh-cn': 'chi_sim',  # Simplified Chinese
    'zh-tw': 'chi_tra',  # Traditional Chinese
    'ja': 'jpn',    # Japanese
    'ko': 'kor',    # Korean
    'ar': 'ara',    # Arabic
    'he': 'heb',    # Hebrew
    'fa': 'fas',    # Persian (Farsi)
    'hi': 'hin',    # Hindi
    'th': 'tha',    # Thai
    'vi': 'vie',    # Vietnamese
    'nl': 'nld',    # Dutch
    'tr': 'tur',    # Turkish
    'pl': 'pol',    # Polish
    'uk': 'ukr',    # Ukrainian
    'ro': 'ron',    # Romanian
    'bg': 'bul',    # Bulgarian
    'el': 'ell',    # Greek
    'ur': 'urd',    # Urdu
    # Add more languages as needed
}



#################### for pdf

def extract_text_pymupdf(
    file_path: str,
    codec: str = 'utf-8'
) -> Optional[str]:
    """
    Extracts full text from a PDF file using PyMuPDF.
    Each page's text is formatted as a single line in the specified format.

    Parameters:
        file_path (str): The path to the PDF file.
        codec (str): The codec to use for text encoding. Defaults to 'utf-8'.

    Returns:
        Optional[str]: The concatenated formatted text if successful, else None.
    """
    try:
        # Ensure the file exists
        if not os.path.isfile(file_path):
            print(f"The file '{file_path}' does not exist.")
            return None

        # Extract the title from the file name
        title = os.path.basename(file_path)

        # Open the PDF file
        doc = fitz.open(file_path)
        formatted_text = ""

        # Iterate through each page and extract text
        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            page_text = page.get_text("text")
            # Replace newline characters with spaces to make the page text a single line
            single_line_text = ' '.join(page_text.splitlines()).strip()
            formatted_entry = f"Title: {title} | Page: {page_num + 1} | Content: {single_line_text}"
            text_processor.add_text(formatted_entry)
            formatted_text += formatted_entry + "\n"

        if not formatted_text.strip():
            print(f"No text could be extracted from '{file_path}'.")
            return None

        return formatted_text

    except Exception as e:
        print(f"An error occurred while extracting text from PDF: {e}")
        return None



############### for pptx

def extract_text_pptx(
    file_path: str,
    codec: str = 'utf-8'
) -> Optional[str]:
    """
    Extracts text from a PowerPoint (.pptx) file.
    Each slide's text is formatted as a single line in the specified format.

    Parameters:
        file_path (str): The path to the PowerPoint file.
        codec (str): The codec to use for text encoding. Defaults to 'utf-8'.

    Returns:
        Optional[str]: The concatenated formatted text if successful, else None.
    """
    try:
        # Ensure the file exists
        if not os.path.isfile(file_path):
            print(f"The file '{file_path}' does not exist.")
            return None

        # Extract the title from the file name
        title = os.path.basename(file_path)

        # Open the PowerPoint file
        prs = Presentation(file_path)
        formatted_text = ""

        # Iterate through each slide and extract text
        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_text_runs = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text_runs.append(shape.text)
            # Combine all text runs into a single line
            single_line_text = ' '.join(slide_text_runs).strip()
            formatted_entry = f"Title: {title} | Page: {slide_num} | Content: {single_line_text}"
            text_processor.add_text(formatted_entry)
            formatted_text += formatted_entry + "\n"

        if not formatted_text.strip():
            print(f"No text could be extracted from '{file_path}'.")
            return None

        return formatted_text

    except FileNotFoundError:
        print(f"File {file_path} not found.")
    except Exception as e:
        print(f"An error occurred while extracting text from PowerPoint: {e}")
    return None



###################### for docs 

# def extract_text_docx(
#     file_path: str,
#     output_txt: Optional[str] = None,
#     codec: str = 'utf-8'
# ) -> Optional[str]:
#     """
#     Extracts text from a Word (.docx) file and optionally saves it to a text file.
#     Each paragraph is formatted as a single line in the specified format.

#     Parameters:
#         file_path (str): The path to the Word file.
#         output_txt (Optional[str]): The path to save the extracted text as a .txt file. If None, the text is not saved to a file.
#         codec (str): The codec to use for text encoding. Defaults to 'utf-8'.

#     Returns:
#         Optional[str]: The concatenated formatted text if successful, else None.
#     """
#     try:
#         # Ensure the file exists
#         if not os.path.isfile(file_path):
#             print(f"The file '{file_path}' does not exist.")
#             return None

#         # Extract the title from the file name
#         title = os.path.basename(file_path)

#         # Open the Word file
#         doc = Document(file_path)
#         formatted_text = ""

#         # Iterate through each paragraph and extract text
#         for para_num, para in enumerate(doc.paragraphs, start=1):
#             para_text = para.text.strip()
#             if para_text:  # Only include non-empty paragraphs
#                 formatted_entry = f"Title: {title} | Page: {para_num} | Content: {para_text}"
#                 text_processor.add_text(formatted_entry)
#                 formatted_text += formatted_entry + "\n"

#         if not formatted_text.strip():
#             print(f"No text could be extracted from '{file_path}'.")
#             return None

#         # Save the formatted text to a .txt file if output_txt path is provided
#         if output_txt:
#             try:
#                 with open(output_txt, 'w', encoding=codec) as f:
#                     f.write(formatted_text)
#                 print(f"Extracted text saved to '{output_txt}'.")
#             except IOError as e:
#                 print(f"Failed to save text to '{output_txt}': {e}")
#                 return None

#         return formatted_text

#     except FileNotFoundError:
#         print(f"File {file_path} not found.")
#     except Exception as e:
#         print(f"An error occurred while extracting text from Word document: {e}")
#     return None

def extract_text_docx(doc_path):
    doc = Document(doc_path)
    formatted_entries = []
    slide_num = 1
    current_title = ""
    current_content = ""
    # Extract the file name without extension to use as Title
    title = os.path.basename(doc_path)
    for para in doc.paragraphs:
        if para.style.name.startswith('Heading 1'):
            if current_title:
                # Format the current entry and add it to the list
                formatted_entry = f"Title: {title} | Page: {slide_num} | Content: {current_content.strip()}"
                text_processor.add_text(formatted_entry)
                formatted_entries.append(formatted_entry)
                slide_num += 1
                current_content = ""
            current_title = para.text
        else:
            current_content += para.text + " "
    # Add the last entry if exists
    if current_title:
        formatted_entry = f"Title: {current_title} | Page: {slide_num} | Content: {current_content.strip()}"
        formatted_entries.append(formatted_entry)
    # Combine all entries into a single string separated by newline characters
    all_entries = "\n".join(formatted_entries)
    return all_entries


# Supported image formats
image_extensions = ['*.png', '*.jpg', '*.jpeg', '*.tiff', '*.bmp']

# Preprocessing function to enhance OCR accuracy for images
def preprocess_image(image_path):
    """
    Preprocesses an image to improve OCR accuracy.
    Steps:
    - Convert to grayscale
    - Apply median filter for noise reduction
    - Enhance contrast
    - Binarize the image using thresholding
    """
    try:
        image = Image.open(image_path).convert('L')  # Convert to grayscale
        image = image.filter(ImageFilter.MedianFilter())  # Reduce noise
        enhancer = ImageEnhance.Contrast(image)
        image = enhancer.enhance(2)  # Enhance contrast
        image = image.point(lambda x: 0 if x < 140 else 255, '1')  # Binarization
        return image
    except Exception as e:
        print(f"Error preprocessing {image_path}: {e}")
        return None

# Function to extract text from the preprocessed image using OCR with multiple languages
def extract_text_multiple_languages(image, languages=['eng', 'heb']):
    """
    Extracts text from an image using Tesseract OCR with multiple languages.
    """
    try:
        lang_param = '+'.join(languages)
        text = pytesseract.image_to_string(image, lang=lang_param)
        return text.strip()
    except Exception as e:
        print(f"Error during OCR with multiple languages: {e}")
        return ""

# List of languages for OCR
desired_languages = ['eng', 'heb']

# Function to process only image files and format the extracted text
def extract_text_from_image(file_path):
    """
    Extracts text from an image file, formats it, and returns the formatted string.
    """
    _, file_extension = os.path.splitext(file_path)
    file_extension = file_extension.lower()

    if file_extension in ['.jpg', '.jpeg', '.png', '.bmp', '.tiff']:
        # Process image with OCR
        preprocessed_img = preprocess_image(file_path)
        if preprocessed_img is None:
            return None
        
        text = extract_text_multiple_languages(preprocessed_img, languages=desired_languages)
        
        # Format the extracted text
        title = os.path.basename(file_path)
        formatted_text = f"Title: {title} | Content: {text}"
        text_processor.add_text(formatted_text)
        return formatted_text
    else:
        print(f"Unsupported file type: {file_extension}")
        return None
    

# Define desired languages for OCR
desired_languages = ['eng', 'heb', 'spa']  # Add more as needed

def extract_text_from_file(file_path):
    _, file_extension = os.path.splitext(file_path)
    file_extension = file_extension.lower()

    if file_extension == '.pdf':
        text = extract_text_pymupdf(file_path)
        return text, "pdf"

    elif file_extension == '.pptx':
        text = extract_text_pptx(file_path)
        return text, "pptx"

    elif file_extension == '.docx':
        text = extract_text_docx(file_path)
        return text, "docx"

    elif file_extension in ['.jpg', '.jpeg', '.png', '.bmp', '.tiff']:
        # Process image with OCR
        preprocessed_img = preprocess_image(file_path)
        if preprocessed_img is None:
            return "", "preprocessing_failed"
        text = extract_text_from_image(file_path)
        return text, "image"

    else:
        print(f"Unsupported file type: {file_extension}")
        return "", "unsupported"

